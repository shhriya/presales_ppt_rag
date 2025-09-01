# backend/logic/ppt_extractor.py
import os
import zipfile
import cv2
import numpy as np
import pandas as pd
from pptx import Presentation
from ..ocr_utils import extract_table_or_text, clean_text_df

def extract_and_segregate_media(pptx_path, output_base_dir):
    """
    Extracts media from PPTX (ppt/media/*), saves them into categorized folders
    and runs OCR on images (tables vs text).
    Returns a dict mapping image filename -> OCR text/table summary.
    """
    file_types = {
        "images": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"],
        "videos": [".mp4", ".mov", ".avi", ".wmv", ".mkv"],
        "audio": [".mp3", ".wav", ".aac", ".ogg", ".m4a"],
        "vectors": [".svg", ".emf", ".wmf"],
        "documents": [".pdf", ".docx", ".xlsx", ".pptx", ".html", ".htm"],
    }
    for folder in file_types.keys():
        os.makedirs(os.path.join(output_base_dir, folder), exist_ok=True)
    os.makedirs(os.path.join(output_base_dir, "others"), exist_ok=True)

    ocr_results = {}
    with zipfile.ZipFile(pptx_path, 'r') as pptx:
        for file in pptx.namelist():
            if file.startswith('ppt/media/'):
                file_data = pptx.read(file)
                file_name = os.path.basename(file)
                ext = os.path.splitext(file_name)[1].lower()

                category = next((folder for folder, exts in file_types.items() if ext in exts), "others")
                save_path = os.path.join(output_base_dir, category, file_name)
                with open(save_path, 'wb') as f:
                    f.write(file_data)

                if category == "images":
                    try:
                        img = cv2.imdecode(np.frombuffer(file_data, np.uint8), cv2.IMREAD_COLOR)
                        mode, result = extract_table_or_text(img)
                        if mode == "table":
                            table_csv = save_path + ".csv"
                            # write CSV of table
                            result.to_csv(table_csv, index=False, encoding="utf-8-sig")

                            table_text = "\n".join(["\t".join(map(str,row)) for row in result.values])
                            ocr_results[file_name] = f"[Table content]\n: {table_text}"
                        else:
                            ocr_results[file_name] = " ".join(result)
                    except Exception as e:
                        print(f"OCR error on {file_name}: {e}")
    return ocr_results


def extract_text_from_shape(shape):
    """
    Pulls all text from shapes (text boxes) and tables in slides.
    """
    if shape.has_text_frame:
        return shape.text.strip()
    elif shape.has_table:
        rows = ["\t".join([cell.text.strip() for cell in row.cells]) for row in shape.table.rows]
        return "\n".join(rows)
    return ""


def extract_presentation_content(ppt_path, img_texts):
    """
    Extracts text content from a PPTX file (slides) and injects OCR results
    for any images referenced in slides.
    Returns list of slide dicts: {"slide_number": idx, "full_text": "..."}
    """
    prs = Presentation(ppt_path)
    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                slide_content.append(text)
        for rel in slide.part.rels.values():
            if "image" in rel.reltype:
                img_name = os.path.basename(rel.target_ref)
                if img_name in img_texts:
                    slide_content.append(f"[Image OCR/Text] {img_texts[img_name]}")
        slides_data.append({"slide_number": idx, "full_text": "\n".join(slide_content)})
    return slides_data
