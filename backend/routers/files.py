# routers/files.py
import os
import sys
import shutil
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from backend.config import SESSIONS_DIR
from datetime import datetime
import mimetypes
import fitz  # PyMuPDF
from pptx import Presentation
import subprocess
import tempfile
import os

# Define uploads directory for group files
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
UPLOADS_DIR = os.path.join(BASE_DIR, "uploads")
router = APIRouter()

# Only show user-uploaded files inside sessions (ignore slides/chunks)
IGNORE_EXTS = {".json", ".index", ".tmp", ".log"}
IGNORE_FILES = {"chunks.json", "slides.json", "Thumbs.db", ".DS_Store"}


def list_session_files(sessions_dir):
    import json
    files = []
    for sid in os.listdir(sessions_dir):
        session_path = os.path.join(sessions_dir, sid)
        if not os.path.isdir(session_path):
            continue
        
        # Load metadata
        meta_path = os.path.join(session_path, "meta.json")
        meta = {}
        if os.path.exists(meta_path):
            with open(meta_path, "r") as mf:
                for m in json.load(mf):
                    meta[m["filename"]] = m
        
        for fname in os.listdir(session_path):
            fpath = os.path.join(session_path, fname)
            if not os.path.isfile(fpath) or fname in IGNORE_FILES or os.path.splitext(fname)[1].lower() in IGNORE_EXTS:
                continue
            
            meta_entry = meta.get(fname, {})
            uploader_info = meta_entry.get("uploaded_by", {})
            uploaded_at = meta_entry.get("uploaded_at", datetime.fromtimestamp(os.path.getmtime(fpath)).isoformat())
            source = meta_entry.get("source", "chat")

            # Only include files uploaded from chatbot; exclude other sources
            if source != "chat":
                continue
            
            files.append({
                "id": f"{sid}_{fname}",
                "session_id": sid,
                "original_filename": fname,
                "preview_url": f"/files/sessions/{sid}/{fname}/preview",
                "download_url": f"/files/sessions/{sid}/{fname}/download",
                "uploaded_at": uploaded_at,
                "uploaded_by": uploader_info
            })
    return files



@router.get("/api/files")
def api_list_files():
    """Return all user-uploaded files inside sessions, ignore system files."""
    files_list = list_session_files(SESSIONS_DIR)
    return {"files": files_list}

@router.get("/api/files/my")
def api_get_my_files():
    """Get all files the current user has access to"""
    # For now, return all files since we don't have user-specific filtering yet
    files_list = list_session_files(SESSIONS_DIR)
    return files_list

@router.get("/files/{file_id}")
def get_file_by_id(file_id: str):
    """Get file by file_id (supports both session files and group files)"""
    try:
        file_path = _resolve_file_path_from_id(file_id)
        
        # Determine content type
        media_type, _ = mimetypes.guess_type(file_path)
        if media_type is None:
            media_type = "application/octet-stream"
        
        # Stream inline for browser preview (no filename to avoid download)
        return FileResponse(
            file_path,
            media_type=media_type,
            headers={"Content-Disposition": "inline"}
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/files/{file_id}/download")
def download_file_by_id(file_id: str):
    """Download file by file_id (supports both session files and group files)"""
    try:
        file_path = _resolve_file_path_from_id(file_id)
        
        # Extract filename from file_id
        if file_id.startswith("group_"):
            # Group files: group_{uuid}_{filename}
            parts = file_id.split("_", 2)
            filename = parts[2] if len(parts) > 2 else "download"
        else:
            # Session files: session_id_filename
            parts = file_id.split("_", 1)
            filename = parts[1] if len(parts) > 1 else "download"
        
        return FileResponse(
            file_path, 
            media_type="application/octet-stream",
            filename=filename
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading file: {str(e)}")

@router.get("/files/sessions/{session_id}/{filename}/preview")
def preview_session_file(session_id: str, filename: str):
    session_path = os.path.join(SESSIONS_DIR, session_id)
    file_path = os.path.join(session_path, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    media_type, _ = mimetypes.guess_type(file_path)
    if media_type is None:
        media_type = "application/octet-stream"  # fallback
    return FileResponse(file_path, media_type=media_type, headers={"Content-Disposition": "inline"})


def _resolve_file_path_from_id(file_id: str) -> str:
    # Handle group files (format: group_{uuid}_{filename})
    if file_id.startswith("group_"):
        # Group files are stored in uploads directory
        file_path = os.path.join(UPLOADS_DIR, file_id)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="Group file not found")
        return file_path
    
    # Handle session files (format: session_id_filename)
    if "_" not in file_id:
        raise HTTPException(status_code=400, detail="Invalid file_id format")
    session_id, filename = file_id.split("_", 1)
    file_path = os.path.join(SESSIONS_DIR, session_id, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Session file not found")
    return file_path

@router.get("/files/{file_id}/as-pdf")
def get_office_as_pdf(file_id: str):
    """Convert Office docs (PPTX/PPT/DOCX/DOC) to PDF for inline viewing.

    Strategy:
    1) On Windows, try Office COM (PowerPoint/Word) if available.
    2) Try LibreOffice (soffice) headless conversion.
    3) For PPTX/PPT only, fallback to simple text-only PDF.
    """
    try:
        src_path = _resolve_file_path_from_id(file_id)
        ext = os.path.splitext(src_path)[1].lower()
        if ext not in [".pptx", ".ppt", ".docx", ".doc"]:
            # If already PDF, just return the original
            if ext == ".pdf":
                return FileResponse(src_path, media_type="application/pdf")
            # Fallback to original
            media_type, _ = mimetypes.guess_type(src_path)
            if media_type is None:
                media_type = "application/octet-stream"
            return FileResponse(src_path, media_type=media_type)

        # Preferred cache path
        out_pdf = src_path + ".preview.pdf"
        if not os.path.exists(out_pdf):
            converted = False

            # 1) On Windows, try Office COM via pywin32 (most reliable)
            if sys.platform.startswith("win") and not converted:
                try:
                    import win32com.client  # type: ignore
                    ext_lower = ext
                    pdf_temp = out_pdf
                    if ext_lower in [".pptx", ".ppt"]:
                        from win32com.client import constants  # type: ignore
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                        powerpoint.Visible = 1
                        try:
                            presentation = powerpoint.Presentations.Open(src_path, WithWindow=False)
                            try:
                                presentation.SaveAs(pdf_temp, constants.ppSaveAsPDF)
                            finally:
                                presentation.Close()
                        finally:
                            powerpoint.Quit()
                        if os.path.exists(pdf_temp):
                            converted = True
                    elif ext_lower in [".docx", ".doc"]:
                        word = win32com.client.Dispatch("Word.Application")
                        word.Visible = 0
                        try:
                            doc = word.Documents.Open(src_path)
                            try:
                                wdFormatPDF = 17
                                doc.SaveAs(pdf_temp, FileFormat=wdFormatPDF)
                            finally:
                                doc.Close(False)
                        finally:
                            word.Quit()
                        if os.path.exists(pdf_temp):
                            converted = True
                except Exception:
                    converted = False

            # 1b) If pywin32 not available, try comtypes as secondary on Windows
            if sys.platform.startswith("win") and not converted:
                try:
                    import comtypes.client  # type: ignore
                    pp_save_as_pdf = 32
                    pdf_temp = out_pdf
                    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                    powerpoint.Visible = 1
                    try:
                        presentation = powerpoint.Presentations.Open(src_path, WithWindow=False)
                        try:
                            presentation.SaveAs(pdf_temp, pp_save_as_pdf)
                        finally:
                            presentation.Close()
                    finally:
                        powerpoint.Quit()
                    if os.path.exists(pdf_temp):
                        converted = True
                except Exception:
                    converted = False

            # 2) Try LibreOffice if available (cross-platform)
            if not converted:
                try:
                    # Resolve soffice executable from env or PATH, accepting a directory path too
                    soffice_env = os.environ.get("LIBREOFFICE_PATH")
                    soffice_candidates = []
                    if soffice_env:
                        # If a directory was supplied, append common executable names
                        if os.path.isdir(soffice_env):
                            soffice_candidates.extend([
                                os.path.join(soffice_env, "soffice"),
                                os.path.join(soffice_env, "soffice.exe"),
                                os.path.join(soffice_env, "program", "soffice"),
                                os.path.join(soffice_env, "program", "soffice.exe"),
                            ])
                        else:
                            soffice_candidates.append(soffice_env)
                    # Also try PATH
                    which_soffice = shutil.which("soffice")
                    if which_soffice:
                        soffice_candidates.append(which_soffice)
                    # Common Windows install location
                    if sys.platform.startswith("win"):
                        soffice_candidates.extend([
                            r"C:\\Program Files\\LibreOffice\\program\\soffice.exe",
                            r"C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
                        ])

                    soffice_exec = next((p for p in soffice_candidates if p and os.path.exists(p)), None)
                    if not soffice_exec:
                        raise RuntimeError("LibreOffice (soffice) not found. Set LIBREOFFICE_PATH to the install folder or soffice.exe.")

                    # Use component export filters; fallback to generic pdf if filter not needed
                    # --norestore/--nolockcheck avoid dialogs; --headless makes it non-interactive
                    proc = subprocess.run([
                        soffice_exec,
                        "--headless",
                        "--norestore",
                        "--nolockcheck",
                        "--convert-to", "pdf",
                        "--outdir", os.path.dirname(src_path),
                        src_path,
                    ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

                    if proc.returncode != 0:
                        raise RuntimeError(f"LibreOffice failed: {proc.stderr or proc.stdout}")

                    # LibreOffice writes <basename>.pdf in the outdir
                    candidate = os.path.splitext(src_path)[0] + ".pdf"
                    if not os.path.exists(candidate):
                        # Some versions may output with uppercase/lowercase variations
                        alt = os.path.splitext(os.path.basename(src_path))[0] + ".pdf"
                        alt_path = os.path.join(os.path.dirname(src_path), alt)
                        if os.path.exists(alt_path):
                            candidate = alt_path

                    if os.path.exists(candidate):
                        os.replace(candidate, out_pdf)
                        converted = True
                    else:
                        raise RuntimeError("LibreOffice did not produce a PDF output file")
                except Exception as e:
                    # Leave a breadcrumb in stderr to help diagnose
                    print(f"[as-pdf] LibreOffice conversion error: {e}")
                    converted = False

            # 3) Fallback: simple text-only conversion (last resort, PPT only)
            if not converted and ext in [".pptx", ".ppt"]:
                try:
                    prs = Presentation(src_path)
                    doc = fitz.open()
                    for slide in prs.slides:
                        page = doc.new_page()
                        y = 72
                        for shape in slide.shapes:
                            try:
                                text = ""
                                if getattr(shape, "has_text_frame", False):
                                    text = shape.text
                                elif getattr(shape, "has_table", False):
                                    rows = []
                                    for row in shape.table.rows:
                                        rows.append(" \t ".join([cell.text for cell in row.cells]))
                                    text = "\n".join(rows)
                                if text:
                                    page.insert_text((72, y), text, fontsize=12)
                                    y += 18 * (text.count("\n") + 1)
                                    if y > page.rect.height - 72:
                                        page = doc.new_page()
                                        y = 72
                            except Exception:
                                continue
                    doc.save(out_pdf)
                    doc.close()
                    converted = True
                except Exception:
                    converted = False

            if not converted:
                raise HTTPException(status_code=500, detail="Unable to convert PPT to PDF on this server")

        return FileResponse(out_pdf, media_type="application/pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {e}")

# routers/files.py
from fastapi import UploadFile, Depends, Request

# Example of getting user info from headers (authHeaders from frontend)
def get_user_from_headers(request: Request):
    user_id = request.headers.get("X-User-Id")
    role = request.headers.get("X-User-Role")
    username = request.headers.get("X-User-Name")  # optional, send from frontend
    email = request.headers.get("X-User-Email")
    return {
        "user_id": user_id,
        "role": role,
        "username": username,
        "email": email
    }

@router.post("/upload-file")
async def upload_file(file: UploadFile, request: Request):
    user = get_user_from_headers(request)
    
    session_id = generate_session_id()  # however you create session folders
    session_path = os.path.join(SESSIONS_DIR, session_id)
    os.makedirs(session_path, exist_ok=True)
    
    file_path = os.path.join(session_path, file.filename)
    
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    # Save metadata (including uploader info)
    meta_path = os.path.join(session_path, "meta.json")
    metadata = []
    if os.path.exists(meta_path):
        import json
        with open(meta_path, "r") as mf:
            metadata = json.load(mf)
    
    metadata.append({
        "filename": file.filename,
        "uploaded_at": datetime.now().isoformat(),
        "uploaded_by": user
    })

    with open(meta_path, "w") as mf:
        import json
        json.dump(metadata, mf)

    return {"detail": "File uploaded", "file_id": f"{session_id}_{file.filename}"}
