# backend/ppt_logic.py
from pptx import Presentation
import warnings, logging, json, os, zipfile, re
import cv2, numpy as np, pandas as pd
from PIL import Image
import pytesseract, easyocr, faiss
from dotenv import load_dotenv
from openai import OpenAI
 
# ===== ENV =====
load_dotenv()
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
os.environ["TF_CPP_MIN_LOG_LEVEL"] = "3"
warnings.filterwarnings("ignore", category=UserWarning)
logging.getLogger("PIL").setLevel(logging.ERROR)
 

# ===== CONFIG =====
CHUNK_SIZE = 500
CHUNK_OVERLAP = 75
 
# ===== OCR INIT =====
reader = easyocr.Reader(['en'], gpu=False)
tess_cmd = os.getenv("TESSERACT_CMD")
if tess_cmd and os.path.exists(tess_cmd):
    pytesseract.pytesseract.tesseract_cmd = tess_cmd
 
# ===== OpenAI client =====
client = OpenAI()
 
# ------------------ UTIL FUNCTIONS ------------------ #
def cluster_positions(pos_list, thresh=10): #Groups detected rows/columns when processing tables from images.
    pos_list = sorted(pos_list)
    clustered = [pos_list[0]]
    for p in pos_list[1:]:
        if abs(p - clustered[-1]) > thresh:
            clustered.append(p)
    return clustered
 
def clean_text_df(df): #Cleans OCR output from tables (removes noise, strips whitespace).
    df = df.replace(r"[|\]\[\}\{I™]+", "", regex=True)
    df = df.applymap(lambda x: x.strip() if isinstance(x, str) else x)
    df = df.dropna(how="all", axis=0).dropna(how="all", axis=1)
    return df.reset_index(drop=True)
 
def extract_table_or_text(img, force_text=False, force_table=False): #Decides whether an image contains a table or just text, extracts accordingly.
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) if len(img.shape) == 3 else img
    if force_text:
        text = pytesseract.image_to_string(gray, config="--psm 3").strip()
        lines = [line for line in text.split("\n") if line.strip()]
        return "text", lines
 
    # Detect lines for table
    thresh = cv2.adaptiveThreshold(~gray, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, 15, -2)
    horizontal = thresh.copy()
    vertical = thresh.copy()
    cols, rows = horizontal.shape[1], vertical.shape[0]
    horizontalStructure = cv2.getStructuringElement(cv2.MORPH_RECT, (max(1, cols // 30), 1))
    verticalStructure = cv2.getStructuringElement(cv2.MORPH_RECT, (1, max(1, rows // 30)))
    horizontal = cv2.dilate(cv2.erode(horizontal, horizontalStructure), horizontalStructure)
    vertical = cv2.dilate(cv2.erode(vertical, verticalStructure), verticalStructure)
    mask = horizontal + vertical
 
    contours, _ = cv2.findContours(mask, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)
    boxes = [(x, y, w, h) for c in contours for (x, y, w, h) in [cv2.boundingRect(c)] if w > 40 and h > 20]
 
    if force_table or len(boxes) >= 5:
        row_positions = cluster_positions([y for (_, y, _, h) in boxes] + [y+h for (_, y, _, h) in boxes])
        col_positions = cluster_positions([x for (x, _, w, _) in boxes] + [x+w for (x, _, w, _) in boxes])
        table = []
        for i in range(len(row_positions)-1):
            row = []
            for j in range(len(col_positions)-1):
                x1, y1, x2, y2 = col_positions[j], row_positions[i], col_positions[j+1], row_positions[i+1]
                roi = gray[y1:y2, x1:x2]
                if roi.size == 0:
                    row.append("")
                    continue
                text = pytesseract.image_to_string(roi, config="--psm 6").strip()
                row.append(text if text else "")
            table.append(row)
        df = pd.DataFrame(table)
        return "table", clean_text_df(df)
 
    text = pytesseract.image_to_string(gray, config="--psm 3").strip()
    lines = [line for line in text.split("\n") if line.strip()]
    return "text", lines
 
# ------------------ 1. Extract media ------------------ #
def extract_and_segregate_media(pptx_path, output_base_dir):
    file_types = {
        "images": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"],
        "videos": [".mp4", ".mov", ".avi", ".wmv", ".mkv"],
        "audio": [".mp3", ".wav", ".aac", ".ogg", ".m4a"],
        "vectors": [".svg", ".emf", ".wmf"],
        "documents": [".pdf", ".docx", ".xlsx", ".pptx", ".html", ".htm"],
    }
    for folder in file_types.keys():
        os.makedirs(os.path.join(output_base_dir, folder), exist_ok=True)
    os.makedirs(os.path.join(output_base_dir, "others"), exist_ok=True)

    ocr_results = {}
    with zipfile.ZipFile(pptx_path, 'r') as pptx:
        for file in pptx.namelist():
            if file.startswith('ppt/media/'):
                file_data = pptx.read(file)
                file_name = os.path.basename(file)
                ext = os.path.splitext(file_name)[1].lower()

                category = next((folder for folder, exts in file_types.items() if ext in exts), "others")
                save_path = os.path.join(output_base_dir, category, file_name)
                with open(save_path, 'wb') as f:
                    f.write(file_data)
 
                if category == "images":
                    try:
                        img = cv2.imdecode(np.frombuffer(file_data, np.uint8), cv2.IMREAD_COLOR)
                        mode, result = extract_table_or_text(img)
                        if mode == "table":
                            table_csv = save_path + ".csv"
                            result.to_csv(table_csv, index=False, encoding="utf-8-sig")

                            table_text = "\n".join(["\t".join(map(str,row)) for row in result.values])
                            ocr_results[file_name] = f"[Table content]\n: {table_text}"
                        else:
                            ocr_results[file_name] = " ".join(result)
                    except Exception as e:
                        print(f"OCR error on {file_name}: {e}")
    return ocr_results
 
# ------------------ 2. Extract PPT text ------------------ #
def extract_text_from_shape(shape): #Pulls all text from shapes (text boxes) and tables in slides.
    if shape.has_text_frame:
        return shape.text.strip()
    elif shape.has_table:
        rows = ["\t".join([cell.text.strip() for cell in row.cells]) for row in shape.table.rows]
        return "\n".join(rows)
    return ""
 
def extract_presentation_content(ppt_path, img_texts):
    prs = Presentation(ppt_path)
    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                slide_content.append(text)
        for rel in slide.part.rels.values():
            if "image" in rel.reltype:
                img_name = os.path.basename(rel.target_ref)
                if img_name in img_texts:
                    slide_content.append(f"[Image OCR/Text] {img_texts[img_name]}")
        slides_data.append({"slide_number": idx, "full_text": "\n".join(slide_content)})
    return slides_data
 
# ------------------ 3. Chunk text ------------------ #
def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    chunks, start = [], 0
    
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap
    
    return chunks
 
# ------------------ 4. Build FAISS ------------------ #
def build_faiss_index(slides, index_path, chunks_json_path):
    texts, metadata = [], []
    for slide in slides:
        chunks = chunk_text(slide["full_text"])
        for chunk in chunks:
            texts.append(chunk)
            metadata.append({"slide": slide["slide_number"]})
    vectors = []
    for text in texts:
        emb = client.embeddings.create(model="text-embedding-3-small", input=text)
        vectors.append(emb.data[0].embedding)
    vectors = np.array(vectors).astype("float32")
    index = faiss.IndexFlatL2(vectors.shape[1])
    index.add(vectors)
    faiss.write_index(index, index_path)
    with open(chunks_json_path, "w", encoding="utf-8") as f:
        json.dump({"texts": texts, "metadata": metadata}, f, ensure_ascii=False, indent=4)
    return index, texts, metadata
 
# ------------------ 5. Search & Answer ------------------ #
def search_and_answer(query, index, texts, metadata):
    print("RAW USER QUERY:", query)

    # Extract just the new question
    match_new = re.search(r"New question:(.*)", query, re.IGNORECASE | re.DOTALL)
    actual_question = match_new.group(1).strip() if match_new else query

    # --- retrieval ---
    match = re.search(r"slide\s+(\d+)", actual_question.lower())
    if match:
        slide_num = int(match.group(1))
        retrieved = [texts[i] for i, meta in enumerate(metadata) if meta["slide"] == slide_num]
        if not retrieved:
            retrieved = texts[:3]
    else:
        emb = client.embeddings.create(model="text-embedding-3-small", input=actual_question).data[0].embedding
        print("EMBEDDING CREATED FOR:", actual_question)
        D, I = index.search(np.array([emb], dtype="float32"), 3)
        retrieved = [texts[i] for i in I[0]]

    prompt = f"""
You are an assistant helping the user understand their PowerPoint presentation.
If the PPT content is messy, incomplete, or seems mismatched, you may reason carefully and infer the intended meaning.

PPT Content:
{retrieved}

Conversation context (last Q&A + new question):
{query}

Guidelines:
1. If the new question is a follow-up, treat it as a request to expand or clarify your **last answer**.
2. Use the retrieved PPT content to elaborate, give examples, or simplify further.
3. If the question can be answered from PPT, do it clearly and concisely and with careful reasoning.
4. Keep tone factual, simple, and user-friendly.

Answer:
"""
    print("FINAL PROMPT TO LLM:", prompt)

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}]
    )
    print("RAW GPT RESPONSE:", response.choices[0].message.content)
    return response.choices[0].message.content
